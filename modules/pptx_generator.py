"""
PPTX Generator — Creates 1-slider slides from article data.

Uses the slide_template.pptx (converted from .potx) as the base.
All element positions are fixed from the template. Text is handled
intelligently: font size fallback, conditional bullets, word-limit
awareness, and normAutofit for PowerPoint-level shrinking.

Shape reference (from template):
  Titre 42        (id=43)  — Article headline
  Rectangle 5     (id=6)   — Category tag (red box, top-left)
  ZoneTexte 6     (id=7)   — Date (below category)
  Rectangle 7     (id=8)   — "SUMMARY & RELEVANT INFORMATION" label
  Rectangle 8     (id=9)   — "IMPLICATIONS" label
  Rectangle 9     (id=10)  — "CREDIBILITY" label
  Rectangle 1     (id=2)   — "RELEVANCE" label
  Rectangle 13    (id=14)  — Summary + relevant info content
  Rectangle 16    (id=17)  — Implications content
  ZoneTexte 4     (id=15)  — Source URL
  Picture 111      (id=16)  — Article image (single, left side)
  Star: 5 Points 10/11/12  — Credibility stars
  Star: 5 Points 22/24/25  — Relevance stars
  TGR Logo                  — Inherited from slide layout (no action needed)
"""

import logging
import os
from datetime import datetime
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional
from urllib.parse import urlparse

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.oxml.ns import qn
from lxml import etree

from config import TemplateLayout as TL, Grid, TEMPLATE_PATH, SLIDES_DIR

logger = logging.getLogger(__name__)


# ── Data Model ──────────────────────────────────────────────────────────

@dataclass
class ArticleData:
    title: str = ""
    summary: str = ""
    relevant_info: str = ""
    implications: str = ""
    implications_sub: list = field(default_factory=list)
    category: str = "GENERAL INNOVATION"
    publication_date: str = ""
    source_url: str = ""
    source_name: str = ""
    credibility_score: float = 3.0
    relevancy_score: float = 3.0
    article_image: Optional[str] = None


# ── Error text filtering ────────────────────────────────────────────────

_ERROR_INDICATORS = [
    "parsing failed", "evaluation parsing failed", "using default scores",
    "failed to", "error:", "exception:", "traceback",
    "no data available", "could not parse", "api error",
]


def _is_error_text(text: str) -> bool:
    if not text:
        return False
    low = text.lower().strip()
    return any(ind in low for ind in _ERROR_INDICATORS)


def _sanitize(article: ArticleData) -> ArticleData:
    if _is_error_text(article.relevant_info):
        logger.warning(f"Filtered error from relevant_info: {article.relevant_info[:60]}")
        article.relevant_info = ""
    if article.summary:
        lines = article.summary.split("\n")
        article.summary = "\n".join(l for l in lines if not _is_error_text(l))
    if _is_error_text(article.implications):
        logger.warning(f"Filtered error from implications: {article.implications[:60]}")
        article.implications = ""
    if article.implications_sub:
        article.implications_sub = [s for s in article.implications_sub if not _is_error_text(s)]
    return article


# ── Helpers ─────────────────────────────────────────────────────────────

def _format_date(date_str: str) -> str:
    if not date_str:
        return datetime.now().strftime("%d/%m/%Y")
    if len(date_str) == 10 and date_str[2] == "/" and date_str[5] == "/":
        return date_str
    for fmt in ["%Y-%m-%d", "%Y-%m-%dT%H:%M:%S", "%Y-%m-%dT%H:%M:%SZ",
                "%d-%m-%Y", "%m/%d/%Y"]:
        try:
            dt = datetime.strptime(date_str.split("+")[0].split("Z")[0], fmt)
            return dt.strftime("%d/%m/%Y")
        except ValueError:
            continue
    return date_str



def _find_shape_by_name(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _find_shape_by_id(slide, shape_id: int):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    return None


def score_to_stars(score: float, max_stars: int = 3) -> int:
    """0-5 score → 1-3 stars."""
    if score <= 1.7:
        return 1
    elif score <= 3.3:
        return 2
    else:
        return 3


# ── Font size fallback ──────────────────────────────────────────────────

# Character capacity for the content boxes (Rectangle 13 & 16) at each font size.
# Box inner area: ~5.58" wide × 2.37" high (5107887 × 2168391 EMUs).
# Derived from: lines_per_box × avg_chars_per_line at each pt size.
#   14pt → line-height ~0.233" → ~10 lines × ~62 chars  = 620
#   12pt → line-height ~0.200" → ~11 lines × ~75 chars  = 825
#   11pt → line-height ~0.183" → ~12 lines × ~82 chars  = 984
#   10pt → line-height ~0.167" → ~14 lines × ~90 chars  = 1260
_CHAR_CAPACITY = {
    1400: 620,
    1200: 825,
    1100: 984,
    1000: 1260,
}


def _choose_font_size(text: str, default_size: int = 1400, min_size: int = 1000) -> int:
    """Pick the largest font size where the text fits within the content box.

    Uses character count as the sizing signal — more accurate than word count
    since it accounts for long words and varied text density.
    Returns font size in hundredths of a point (e.g. 1400 = 14pt).
    """
    chars = len(text.strip())
    for size in sorted(_CHAR_CAPACITY.keys(), reverse=True):
        if size < min_size or size > default_size:
            continue
        if chars <= _CHAR_CAPACITY[size]:
            return size
    return min_size


# ── XML builders ────────────────────────────────────────────────────────

def _enable_autofit(bodyPr, text: str = "", capacity: int = 0):
    """Enable normAutofit on bodyPr so PowerPoint shrinks text to fit at file-open.

    PowerPoint only applies normAutofit shrink at load time when ``fontScale``
    is set explicitly on the element. Without it the scale stays at 100% until
    the user manually edits the text. This function pre-computes ``fontScale``
    from the ratio of actual character count to the box capacity at the chosen
    font size, using a 90% safety margin so slightly-optimistic estimates still
    produce a result that fits.

    Args:
        bodyPr:   The <a:bodyPr> element on the shape.
        text:     The full text being written into the shape.
        capacity: Character capacity of the box at the chosen font size
                  (from _CHAR_CAPACITY).
    """
    if bodyPr is None:
        return

    # Remove any existing autofit elements at any nesting depth
    autofit_tags = {qn(t) for t in ("a:noAutofit", "a:normAutofit", "a:spAutoFit")}
    for el in list(bodyPr.iter()):
        if el.tag in autofit_tags and el is not bodyPr:
            el.getparent().remove(el)

    autofit = etree.Element(qn("a:normAutofit"))

    chars = len(text.strip())
    if chars > 0 and capacity > 0:
        # Trigger fontScale when chars exceed 90% of capacity — this catches
        # cases where our capacity estimate is slightly optimistic and ensures
        # the text never clips the box edge.
        effective = capacity * 0.9
        if chars > effective:
            scale = int(100000 * (effective / chars))
            scale = max(scale, 62500)          # PowerPoint minimum = 62.5%
            scale = (scale // 2500) * 2500     # PowerPoint steps in 2.5%
            autofit.set("fontScale", str(scale))
            autofit.set("lnSpcReduction", "10000")  # 10% line-spacing reduction

    bodyPr.insert(0, autofit)


def _make_run(text: str, font_size: int, lang: str = "en-US",
              bold: bool = False) -> etree._Element:
    """Create an <a:r> run element."""
    r = etree.Element(qn("a:r"))
    rPr = etree.SubElement(r, qn("a:rPr"))
    rPr.set("lang", lang)
    rPr.set("sz", str(font_size))
    rPr.set("dirty", "0")
    if bold:
        rPr.set("b", "1")
    solidFill = etree.SubElement(rPr, qn("a:solidFill"))
    schemeClr = etree.SubElement(solidFill, qn("a:schemeClr"))
    schemeClr.set("val", "tx1")
    t = etree.SubElement(r, qn("a:t"))
    t.text = text
    return r


def _make_paragraph(text: str, font_size: int, lang: str = "en-US",
                    space_before: int = 400) -> etree._Element:
    """Create a plain paragraph (no bullet)."""
    p = etree.Element(qn("a:p"))
    pPr = etree.SubElement(p, qn("a:pPr"))
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
    spcPts.set("val", str(space_before))
    p.append(_make_run(text, font_size, lang))
    return p


def _make_bullet_paragraph(text: str, font_size: int, level: int = 0,
                           lang: str = "en-GB") -> etree._Element:
    """Create a bulleted paragraph. level=0 uses •, level=1 uses Wingdings Ø."""
    p = etree.Element(qn("a:p"))
    pPr = etree.SubElement(p, qn("a:pPr"))

    if level == 0:
        pPr.set("marL", "182563")
        pPr.set("indent", "-182563")
    else:
        pPr.set("marL", "446088")
        pPr.set("lvl", "1")
        pPr.set("indent", "-271463")

    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
    spcPts.set("val", "600")

    buFont = etree.SubElement(pPr, qn("a:buFont"))
    if level == 0:
        buFont.set("typeface", "Arial")
        buFont.set("panose", "020B0604020202020204")
        buFont.set("pitchFamily", "34")
        buFont.set("charset", "0")
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "\u2022")
    else:
        buFont.set("typeface", "Wingdings")
        buFont.set("panose", "05000000000000000000")
        buFont.set("pitchFamily", "2")
        buFont.set("charset", "2")
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "\u00D8")

    r = _make_run(text, font_size, lang)
    # Add altLang for sub-bullets (matches template)
    if level == 1:
        rPr = r.find(qn("a:rPr"))
        rPr.set("altLang", "ja-JP")
    p.append(r)
    return p


# ── Shape updaters ──────────────────────────────────────────────────────

def _update_title(slide, title: str):
    """Element 11: headline text in Titre 42."""
    shape = _find_shape_by_name(slide, "Titre 42")
    if not shape or not shape.has_text_frame:
        logger.warning("Could not find Titre 42")
        return

    # Truncate if excessively long
    if len(title) > TL.TITLE_MAX_CHARS:
        title = title[:TL.TITLE_MAX_CHARS - 3] + "..."

    txBody = shape._element.find(qn("p:txBody"))
    bodyPr = txBody.find(qn("a:bodyPr"))
    # Title is short; autofit without pre-computed fontScale is fine
    _enable_autofit(bodyPr)

    # Replace text in existing paragraph structure to preserve formatting
    for p in txBody.findall(qn("a:p")):
        for r in p.findall(qn("a:r")):
            t = r.find(qn("a:t"))
            if t is not None:
                t.text = title
                return
    # Fallback: build from scratch
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)
    p = etree.SubElement(txBody, qn("a:p"))
    r = _make_run(title, TL.TITLE_FONT_SIZE, bold=True)
    rPr = r.find(qn("a:rPr"))
    latin = etree.SubElement(rPr, qn("a:latin"))
    latin.set("typeface", "+mn-lt")
    p.append(r)


def _update_category(slide, category: str):
    """Element 2: category tag in Rectangle 5."""
    cat = category.upper().strip()
    if len(cat) > TL.CATEGORY_MAX_CHARS:
        cat = cat[:TL.CATEGORY_MAX_CHARS - 1] + "."

    shape = _find_shape_by_name(slide, "Rectangle 5")
    if not shape or not shape.has_text_frame:
        logger.warning("Could not find Rectangle 5")
        return

    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.text = cat
            return


def _update_summary(slide, summary: str, relevant_info: str):
    """Element 9: summary + relevant info in Rectangle 13.

    Text is plain paragraphs by default. No bullets unless the text
    explicitly contains sub-points (lines starting with "- ").
    Font size is selected based on total word count.
    """
    shape = _find_shape_by_name(slide, "Rectangle 13")
    if not shape:
        logger.warning("Could not find Rectangle 13")
        return

    txBody = shape._element.find(qn("p:txBody"))
    if txBody is None:
        return

    bodyPr = txBody.find(qn("a:bodyPr"))

    # Remove existing paragraphs
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    # Combine all text for font size calculation
    all_text = summary
    if relevant_info and relevant_info.strip() and not _is_error_text(relevant_info):
        all_text += " " + relevant_info
    font_size = _choose_font_size(all_text, TL.SUMMARY_FONT_SIZE, TL.SUMMARY_FONT_MIN)
    _enable_autofit(bodyPr, text=all_text, capacity=_CHAR_CAPACITY.get(font_size, 620))

    # Build summary paragraphs
    summary_lines = [l.strip() for l in summary.split("\n") if l.strip()]
    has_sub_points = any(l.startswith("- ") for l in summary_lines)

    if has_sub_points:
        for line in summary_lines:
            if line.startswith("- "):
                txBody.append(_make_bullet_paragraph(line[2:], font_size, level=0))
            else:
                txBody.append(_make_paragraph(line, font_size))
    else:
        for line in summary_lines:
            txBody.append(_make_paragraph(line, font_size))

    # Add relevant info if present
    if relevant_info and relevant_info.strip() and not _is_error_text(relevant_info):
        rel_lines = [l.strip() for l in relevant_info.split("\n") if l.strip()]
        has_rel_sub = any(l.startswith("- ") for l in rel_lines)

        if has_rel_sub:
            for line in rel_lines:
                if line.startswith("- "):
                    txBody.append(_make_bullet_paragraph(line[2:], font_size, level=0))
                else:
                    txBody.append(_make_paragraph(line, font_size))
        else:
            for line in rel_lines:
                txBody.append(_make_paragraph(line, font_size))


def _update_implications(slide, main_point: str, sub_points: list):
    """Element 10: implications in Rectangle 16.

    Main point is a plain paragraph. Sub-points get Wingdings Ø bullets
    ONLY if they exist. Font size is selected based on total word count.
    """
    shape = _find_shape_by_name(slide, "Rectangle 16")
    if not shape:
        logger.warning("Could not find Rectangle 16")
        return

    txBody = shape._element.find(qn("p:txBody"))
    if txBody is None:
        return

    bodyPr = txBody.find(qn("a:bodyPr"))

    # Remove existing paragraphs
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    # Calculate font size from total text.
    # Unlike prose, bullet sub-points each occupy their own line regardless of
    # length, plus carry 6pt spcBef spacing before them. Raw character count
    # alone therefore underestimates the vertical space required. We add a
    # per-sub-point overhead (≈1.5 visual lines at typical font sizes) so the
    # font-size picker accounts for the line-break and spacing overhead of the
    # bulleted structure. The real text (without padding) is still passed to
    # _enable_autofit so PowerPoint's fontScale hint remains accurate.
    all_text = main_point + " " + " ".join(sub_points)
    _BULLET_OVERHEAD_CHARS = 130  # ~1.5 visual lines × ~90 chars/line at 10pt
    sizing_text = all_text.strip() + ("x" * (len(sub_points) * _BULLET_OVERHEAD_CHARS))
    font_size = _choose_font_size(sizing_text, TL.IMPLICATIONS_FONT_SIZE, TL.IMPLICATIONS_FONT_MIN)
    _enable_autofit(bodyPr, text=all_text, capacity=_CHAR_CAPACITY.get(font_size, 620))

    # Main point — plain paragraph (no bullet)
    txBody.append(_make_paragraph(main_point, font_size, lang="en-GB", space_before=400))

    # Sub-points — Wingdings Ø bullets ONLY if sub_points exist
    for sub in sub_points:
        txBody.append(_make_bullet_paragraph(sub, font_size, level=1, lang="en-GB"))


def _extract_domain(url: str) -> str:
    """Extract display domain from full URL.
    'https://www.youtube.com/watch?v=abc' → 'https://www.youtube.com/'
    """
    try:
        parsed = urlparse(url)
        return f"{parsed.scheme}://{parsed.netloc}/"
    except Exception:
        return url[:30] + "..." if len(url) > 30 else url


def _update_footer(slide, source_url: str, date_str: str):
    """Build footer: 'Source: [domain] Date: [DD/MM/YYYY]' in ZoneTexte 4.
    Domain is hyperlinked to full source_url. No explicit font size — inherits template default.
    """
    shape = _find_shape_by_name(slide, "ZoneTexte 4")
    if not shape:
        logger.warning("Could not find ZoneTexte 4")
        return

    txBody = shape._element.find(qn("p:txBody"))
    if txBody is None:
        return

    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    p = etree.SubElement(txBody, qn("a:p"))

    # Run 1: "Source: "
    r1 = etree.SubElement(p, qn("a:r"))
    rPr1 = etree.SubElement(r1, qn("a:rPr"))
    rPr1.set("lang", "fr-FR")
    rPr1.set("dirty", "0")
    t1 = etree.SubElement(r1, qn("a:t"))
    t1.text = "Source: "

    # Run 2: domain (hyperlinked to full URL)
    domain = _extract_domain(source_url) if source_url else ""
    r2 = etree.SubElement(p, qn("a:r"))
    rPr2 = etree.SubElement(r2, qn("a:rPr"))
    rPr2.set("lang", "fr-FR")
    rPr2.set("dirty", "0")
    if source_url:
        rel = slide.part.relate_to(
            source_url,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            is_external=True,
        )
        hlinkClick = etree.SubElement(rPr2, qn("a:hlinkClick"))
        hlinkClick.set(qn("r:id"), rel)
    t2 = etree.SubElement(r2, qn("a:t"))
    t2.text = domain

    # Run 3: " Date: "
    r3 = etree.SubElement(p, qn("a:r"))
    rPr3 = etree.SubElement(r3, qn("a:rPr"))
    rPr3.set("lang", "fr-FR")
    rPr3.set("dirty", "0")
    t3 = etree.SubElement(r3, qn("a:t"))
    t3.text = " Date: "

    # Run 4: date value
    r4 = etree.SubElement(p, qn("a:r"))
    rPr4 = etree.SubElement(r4, qn("a:rPr"))
    rPr4.set("lang", "en-US")
    rPr4.set("dirty", "0")
    t4 = etree.SubElement(r4, qn("a:t"))
    t4.text = date_str


def _set_star_color(shape, color_hex: str):
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return
    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is None:
        solidFill = etree.SubElement(spPr, qn("a:solidFill"))
    for child in list(solidFill):
        solidFill.remove(child)
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", color_hex)


def _update_stars(slide, credibility: float, relevance: float):
    """Elements 6 & 8: star ratings."""
    cred_n = score_to_stars(credibility)
    rel_n = score_to_stars(relevance)

    cred_names = ["Star: 5 Points 10", "Star: 5 Points 11", "Star: 5 Points 12"]
    for i, name in enumerate(cred_names):
        shape = _find_shape_by_name(slide, name)
        if shape:
            _set_star_color(shape, TL.STAR_FILLED if (i + 1) <= cred_n else TL.STAR_EMPTY)

    rel_names = ["Star: 5 Points 22", "Star: 5 Points 24", "Star: 5 Points 25"]
    for i, name in enumerate(rel_names):
        shape = _find_shape_by_name(slide, name)
        if shape:
            _set_star_color(shape, TL.STAR_FILLED if (i + 1) <= rel_n else TL.STAR_EMPTY)


# ── Image handling ──────────────────────────────────────────────────────

def _replace_article_image(slide, image_path: str, prs):
    """Replace the single article image (Picture 111) using two-boundary scaling."""
    if not image_path or not os.path.exists(image_path):
        _remove_pic_shape(slide, "Picture 111")
        return

    from PIL import Image as PILImage

    # Find Picture 111 to read its slot position and size
    spTree = slide._element.find(qn("p:cSld")).find(qn("p:spTree"))
    for pic_elem in spTree.findall(qn("p:pic")):
        nvPicPr = pic_elem.find(qn("p:nvPicPr"))
        if nvPicPr is None:
            continue
        cNvPr = nvPicPr.find(qn("p:cNvPr"))
        if cNvPr is None or cNvPr.get("name") != "Picture 111":
            continue

        spPr = pic_elem.find(qn("p:spPr"))
        if spPr is None:
            break
        xfrm = spPr.find(qn("a:xfrm"))
        if xfrm is None:
            break
        off = xfrm.find(qn("a:off"))
        ext = xfrm.find(qn("a:ext"))
        if off is None or ext is None:
            break

        slot_left = int(off.get("x", 0))
        slot_top = int(off.get("y", 0))
        outer_w = int(ext.get("cx", 0))
        outer_h = int(ext.get("cy", 0))
        if outer_w == 0 or outer_h == 0:
            break

        # Remove the placeholder shape
        spTree.remove(pic_elem)

        # Get image pixel dimensions
        try:
            with PILImage.open(image_path) as img:
                img_w, img_h = img.size
        except Exception as e:
            logger.warning(f"Could not read image dimensions: {e}")
            return

        # Inner boundary = 88% of outer (soft minimum)
        inner_w = int(outer_w * 0.88)
        inner_h = int(outer_h * 0.88)

        # Two-boundary scaling
        scale_outer = min(outer_w / img_w, outer_h / img_h)
        scale_inner = max(inner_w / img_w, inner_h / img_h)
        scale = min(scale_outer, max(scale_inner, 1e-6))

        new_w = int(img_w * scale)
        new_h = int(img_h * scale)

        # Center inside outer box
        left_offset = slot_left + int((outer_w - new_w) / 2)
        top_offset = slot_top + int((outer_h - new_h) / 2)

        slide.shapes.add_picture(
            image_path,
            left_offset,
            top_offset,
            width=Emu(new_w),
            height=Emu(new_h),
        )
        logger.info("Replaced article image (two-boundary scaled)")
        return

    logger.warning("Could not find Picture 111 for image replacement")


def _remove_pic_shape(slide, shape_name: str):
    spTree = slide._element.find(qn("p:cSld")).find(qn("p:spTree"))
    for pic_elem in spTree.findall(qn("p:pic")):
        nvPicPr = pic_elem.find(qn("p:nvPicPr"))
        if nvPicPr is not None:
            cNvPr = nvPicPr.find(qn("p:cNvPr"))
            if cNvPr is not None and cNvPr.get("name") == shape_name:
                spTree.remove(pic_elem)
                logger.info(f"Removed image shape: {shape_name}")
                return


def _expand_image_to_bounds(slide):
    """Post-placement expansion: scale the image up to fill the available bounding box.

    Only expands, never shrinks. Preserves aspect ratio. Centers within the box.
    """
    available_width = TL.IMG_BOX_RIGHT - TL.IMG_BOX_LEFT
    available_height = TL.IMG_BOX_BOTTOM - TL.IMG_BOX_TOP

    # Find the picture shape (shape_type 13 = Picture)
    pic_shape = None
    for shape in slide.shapes:
        if shape.shape_type == 13:
            pic_shape = shape
            break

    if pic_shape is None:
        return

    # Read current position and size via XML
    spPr = pic_shape._element.find(qn("p:spPr"))
    if spPr is None:
        return
    xfrm = spPr.find(qn("a:xfrm"))
    if xfrm is None:
        return
    off = xfrm.find(qn("a:off"))
    ext = xfrm.find(qn("a:ext"))
    if off is None or ext is None:
        return

    current_width = int(ext.get("cx", 0))
    current_height = int(ext.get("cy", 0))
    if current_width == 0 or current_height == 0:
        return

    # Compute scale factor — only expand, never shrink
    scale_w = available_width / current_width
    scale_h = available_height / current_height
    scale = min(scale_w, scale_h)

    if scale <= 1.0:
        return

    new_width = int(current_width * scale)
    new_height = int(current_height * scale)
    new_left = TL.IMG_BOX_LEFT + int((available_width - new_width) / 2)
    new_top = TL.IMG_BOX_TOP + int((available_height - new_height) / 2)

    # Update position and size
    off.set("x", str(new_left))
    off.set("y", str(new_top))
    ext.set("cx", str(new_width))
    ext.set("cy", str(new_height))

    logger.info(f"Expanded image: {current_width}x{current_height} → {new_width}x{new_height}")


# ── Main generation ─────────────────────────────────────────────────────

def generate_slide(article: ArticleData, output_filename: str = None) -> Optional[str]:
    """Generate a single slide from article data.

    Returns path to the generated .pptx file, or None on failure.
    """
    try:
        article = _sanitize(article)
        prs = Presentation(str(TEMPLATE_PATH))

        # Keep only slide 1
        slide_ids = prs.slides._sldIdLst
        while len(slide_ids) > 1:
            rId = slide_ids[-1].get(qn("r:id"))
            prs.part.drop_rel(rId)
            slide_ids.remove(slide_ids[-1])

        slide = prs.slides[0]

        # 0. Clean up phantom/legacy shapes
        for shape in list(slide.shapes):
            if shape.shape_id in (107, 108):
                shape._element.getparent().remove(shape._element)
            elif shape.name == "ZoneTexte 6":
                shape._element.getparent().remove(shape._element)

        # 1. Title (Element 11)
        _update_title(slide, article.title)

        # 2. Category (Element 2)
        _update_category(slide, article.category)

        # 3. Summary + Relevant Info (Element 9)
        _update_summary(slide, article.summary, article.relevant_info)

        # 4. Implications (Element 10) — always called so template placeholder is cleared
        _update_implications(slide, article.implications or "",
                             article.implications_sub or [])

        # 5. Footer (Source + Date in ZoneTexte 4)
        _update_footer(slide, article.source_url, _format_date(article.publication_date))

        # 6. Stars (Elements 6 & 8)
        _update_stars(slide, article.credibility_score, article.relevancy_score)

        # 7. Article image
        _replace_article_image(slide, article.article_image, prs)

        # 8. Expand image to fill available space
        _expand_image_to_bounds(slide)

        # Save
        if not output_filename:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe = "".join(c for c in article.title[:40] if c.isalnum() or c in " -_").strip()
            safe = safe.replace(" ", "_")
            output_filename = f"article_{ts}_{safe}.pptx"

        output_path = SLIDES_DIR / output_filename
        prs.save(str(output_path))
        logger.info(f"Slide saved: {output_path}")
        return str(output_path)

    except Exception as e:
        logger.error(f"Slide generation failed: {e}", exc_info=True)
        return None


def generate_slide_from_notion_data(notion_data: dict,
                                     image_paths: dict = None) -> Optional[str]:
    """Generate a slide from raw Notion data dict.

    Args:
        notion_data: Dict with Notion field values.
        image_paths: Dict with optional key "article" pointing to image path.
                     Also accepts legacy keys "main", "byline", etc.
    """
    images = image_paths or {}

    # Parse implications: first line = main point, rest = sub-bullets
    impl_raw = notion_data.get("implications", "")
    impl_lines = [l.strip() for l in impl_raw.split("\n") if l.strip()]
    impl_lines = [l[1:].lstrip() if l.startswith('-') else l for l in impl_lines]
    main_impl = impl_lines[0] if impl_lines else impl_raw
    sub_impl = impl_lines[1:] if len(impl_lines) > 1 else []

    # Resolve image: prefer "article" key, fall back to "main"
    article_img = images.get("article") or images.get("main")

    article = ArticleData(
        title=notion_data.get("title", "Untitled Article"),
        summary=notion_data.get("summary", ""),
        relevant_info=notion_data.get("relevant_info", ""),
        implications=main_impl,
        implications_sub=sub_impl,
        category=notion_data.get("category", "GENERAL INNOVATION"),
        publication_date=notion_data.get("publication_date", ""),
        source_url=notion_data.get("source_url", ""),
        source_name=notion_data.get("source_name", ""),
        credibility_score=float(notion_data.get("credibility_score", 3.0)),
        relevancy_score=float(notion_data.get("relevancy_score", 3.0)),
        article_image=article_img,
    )

    return generate_slide(article)


# ── Screenshot backup slide ───────────────────────────────────────────────

async def add_screenshot_slide(slide_path: str, article_url: str) -> bool:
    """Append two backup screenshot slides to the existing presentation.

    Slide 2 — viewport-only (1920x1080, horizontal laptop view, fills slide
    edge-to-edge since it matches the 16:9 aspect ratio).

    Slide 3 — full-page scroll capture (entire article). Scaled to fill the
    slide width and centered vertically; if the image is taller than the
    slide it is scaled to fit the height instead.

    Both screenshots are kept in memory — no image files are saved to disk.

    Args:
        slide_path: Path to the existing .pptx file to append to.
        article_url: URL of the article to screenshot.

    Returns:
        True on success, False on failure.
    """
    from io import BytesIO
    from playwright.async_api import async_playwright
    from config import BROWSER_HEADLESS, BROWSER_TIMEOUT

    try:
        async with async_playwright() as p:
            browser = await p.chromium.launch(headless=BROWSER_HEADLESS)
            context = await browser.new_context(
                viewport={"width": 1920, "height": 1080},
            )
            page = await context.new_page()

            logger.info(f"Screenshot slide: navigating to {article_url} (desktop 1920x1080)")
            await page.goto(article_url, wait_until="networkidle",
                            timeout=BROWSER_TIMEOUT)
            await page.wait_for_timeout(2000)

            # Dismiss cookie banners using common selectors
            cookie_selectors = [
                '#onetrust-accept-btn-handler',
                '#CybotCookiebotDialogBodyLevelButtonLevelOptinAllowAll',
                '#CybotCookiebotDialogBodyButtonAccept',
                '#didomi-notice-agree-button',
                '.fc-cta-consent',
                'button:has-text("Accept All")',
                'button:has-text("Accept all")',
                'button:has-text("Accept Cookies")',
                'button:has-text("Accept")',
                'button:has-text("I agree")',
                'button:has-text("Agree")',
                'button:has-text("OK")',
                'button:has-text("Got it")',
                'button:has-text("Allow")',
            ]
            for selector in cookie_selectors:
                try:
                    btn = page.locator(selector).first
                    if await btn.is_visible(timeout=500):
                        await btn.click(timeout=1000)
                        await page.wait_for_timeout(500)
                        break
                except Exception:
                    continue

            # Capture both screenshots in memory
            viewport_bytes = await page.screenshot(full_page=False)
            logger.info("Viewport screenshot captured (1920x1080)")

            fullpage_bytes = await page.screenshot(full_page=True)
            logger.info("Full-page screenshot captured")

            await browser.close()

        # Open existing presentation
        prs = Presentation(slide_path)
        blank_layout = prs.slide_layouts[15]  # "Dark Slide" — 0 placeholders

        slide_w = TL.SLIDE_WIDTH   # EMUs
        slide_h = TL.SLIDE_HEIGHT

        # ── Slide 2: viewport screenshot (16:9 — fills edge-to-edge) ──
        slide2 = prs.slides.add_slide(blank_layout)
        slide2.shapes.add_picture(
            BytesIO(viewport_bytes), 0, 0, slide_w, slide_h
        )

        # ── Slide 3: full-page screenshot (scaled to fit) ─────────────
        slide3 = prs.slides.add_slide(blank_layout)

        from PIL import Image
        with Image.open(BytesIO(fullpage_bytes)) as img:
            img_w, img_h = img.size

        scale = slide_w / img_w
        scaled_h = int(img_h * scale)

        if scaled_h > slide_h:
            # Too tall — scale to fit height, center horizontally
            scale = slide_h / img_h
            final_w = int(img_w * scale)
            final_h = slide_h
            left = (slide_w - final_w) // 2
            top = 0
        else:
            # Fits — fill width, center vertically
            final_w = slide_w
            final_h = scaled_h
            left = 0
            top = (slide_h - final_h) // 2

        slide3.shapes.add_picture(
            BytesIO(fullpage_bytes), left, top, final_w, final_h
        )

        prs.save(slide_path)
        logger.info(f"Screenshot slides (2 & 3) added to {slide_path}")
        return True

    except Exception as e:
        logger.error(f"Failed to add screenshot slides: {e}", exc_info=True)
        return False
